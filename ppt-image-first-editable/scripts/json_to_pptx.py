#!/usr/bin/env python3
"""
json_to_pptx.py — 把 Phase C 的 layout.json / deck.json 渲染成 PPTX

这是 Phase C 的渲染器。模型很简单：每页 = 一张背景 Picture 占满整页 + 若干 TextBox。
所有坐标都用 fraction (0-1)，避免分辨率换算的坑。

Deck JSON schema（多页推荐）:
{
  "deck": {
    "ratio": "16:9",                    # 可选；默认 16:9
    "slide_width_in": 13.333,           # 可选；默认按 ratio 推算
    "slide_height_in": 7.5,             # 可选
    "default_font": "PingFang SC",      # 可选；默认 PingFang SC
    "default_font_size_pt": 18,         # 可选；默认 18pt
    "default_color": "#1a1a1a"          # 可选
  },
  "slides": [
    {
      "id": "01-cover",
      "background": "phaseC/backgrounds/01-cover.png",    # 必填；相对 deck.json 解析
      "text_boxes": [
        {
          "id": "tb-1",                  # 可选；调试用
          "text": "标题\\n第二行",        # 必填；用 \n 表示换行
          "x": 0.08, "y": 0.15,         # 必填；左上角 fraction
          "w": 0.50, "h": 0.20,         # 必填；尺寸 fraction
          "font_family": "PingFang SC",  # 可选；默认 default_font
          "font_size_pt": 36,            # 可选；默认 default_font_size_pt
          "color": "#1a1a1a",            # 可选；十六进制颜色
          "bold": true,                  # 可选；默认 false
          "italic": false,               # 可选；默认 false
          "align": "left",               # 可选；left / center / right；默认 left
          "valign": "top",               # 可选；top / middle / bottom；默认 top
          "line_spacing": 1.2,           # 可选；行距倍数，默认 1.2
          "auto_fit": false              # 可选；自动缩放字号填满文本框；默认 false
        }
      ]
    }
  ]
}

Layout JSON schema（单页，简化版）= 上面 slides 数组里的一个对象，本脚本会自动包装成 deck。

用法:
  # 单页
  python3 scripts/json_to_pptx.py phaseC/01.json -o phaseC/01.pptx

  # 多页（deck.json）
  python3 scripts/json_to_pptx.py phaseC/deck.json -o phaseC/<topic>.pptx

  # 渲染预览图（每页一张 PNG，方便对比 HTML 编辑器里看到的效果）
  python3 scripts/json_to_pptx.py phaseC/deck.json -o phaseC/<topic>.pptx \\
      --preview-dir phaseC/preview

依赖: pip3 install python-pptx pillow
"""

from __future__ import annotations

import argparse
import base64
import io
import json
import sys
from pathlib import Path
from typing import Any
from urllib.parse import unquote, urlparse

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    sys.exit("缺少依赖：pip3 install python-pptx pillow")

try:
    from PIL import Image
except ImportError:
    Image = None  # 预览渲染才需要


# ─────────────────────────────────────────────────────────
# 默认值与字体安全集
# ─────────────────────────────────────────────────────────

DEFAULT_RATIO = "16:9"
DEFAULT_FONT = "PingFang SC"
DEFAULT_FONT_SIZE_PT = 18
DEFAULT_COLOR = "#1a1a1a"
DEFAULT_LINE_SPACING = 1.2

# Phase C 限定的系统安全字体集
# 不在此集合的字体会被警告（仍会写入 PPTX，但跨平台可能掉到 fallback）
SAFE_FONT_SET = {
    "PingFang SC", "PingFang TC", "PingFang HK",
    "Microsoft YaHei", "Microsoft YaHei UI",
    "Hiragino Sans GB", "Heiti SC",
    "Source Han Sans CN", "Source Han Sans SC",  # 思源黑（多数 Mac/Linux 有，Win 需安装）
    "Arial", "Helvetica", "Helvetica Neue",
    "Verdana", "Tahoma", "Calibri",
    "Times New Roman", "Georgia",
    "SimSun", "SimHei", "KaiTi", "FangSong",
    "Courier New", "Consolas", "Menlo", "Monaco",
}

RATIO_PRESETS = {
    "16:9": (13.333, 7.5),
    "4:3":  (10.0,  7.5),
    "3:2":  (10.5,  7.0),
    "1:1":  (7.5,   7.5),
}

ALIGN_MAP = {
    "left":   PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right":  PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

VALIGN_MAP = {
    "top":    MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


# ─────────────────────────────────────────────────────────
# 工具
# ─────────────────────────────────────────────────────────

def hex_to_rgb(s: str) -> RGBColor:
    s = s.strip().lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) != 6:
        raise ValueError(f"颜色格式不对（要 #RRGGBB）: {s}")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def resolve_background_source(ref: str, base_dir: Path) -> Path | io.BytesIO:
    """把背景引用解析为本地路径或内存流。

    支持：
    - 相对 / 绝对本地路径
    - file:// URL
    - data:...;base64,... URL
    """
    if not isinstance(ref, str):
        ref = str(ref)

    if ref.startswith("data:"):
        header, _, payload = ref.partition(",")
        if not payload:
            raise SystemExit("data URL 缺少 payload")
        if ";base64" in header:
            data = base64.b64decode(payload)
        else:
            data = unquote(payload).encode("utf-8")
        return io.BytesIO(data)

    parsed = urlparse(ref)
    if parsed.scheme == "file":
        path = Path(unquote(parsed.path))
        return path

    if parsed.scheme in {"http", "https"}:
        raise SystemExit(f"当前渲染器不支持网络背景 URL: {ref}")

    pp = Path(ref)
    if pp.is_absolute():
        return pp
    return (base_dir / pp).resolve()


def clone_background_source(source: Path | io.BytesIO) -> Path | io.BytesIO:
    if isinstance(source, io.BytesIO):
        return io.BytesIO(source.getvalue())
    return source


def warn(msg: str) -> None:
    print(f"[warn] {msg}", file=sys.stderr)


# ─────────────────────────────────────────────────────────
# 主渲染
# ─────────────────────────────────────────────────────────

def load_deck(json_path: Path) -> tuple[dict, list[dict], Path]:
    """读 JSON，规范成 (deck_cfg, slides_list, base_dir)。base_dir 用来解析背景相对路径。"""
    data = json.loads(json_path.read_text(encoding="utf-8"))
    base_dir = json_path.parent

    # 单页 schema 自动包装
    if "slides" not in data and "text_boxes" in data:
        slides = [data]
        deck_cfg = {}
    else:
        slides = data.get("slides", [])
        deck_cfg = data.get("deck", {})

    if not slides:
        raise SystemExit(f"JSON 里没有 slides: {json_path}")
    return deck_cfg, slides, base_dir


def make_presentation(deck_cfg: dict) -> tuple[Presentation, float, float]:
    """根据 deck 配置创建空 Presentation，返回 (prs, slide_w_in, slide_h_in)。"""
    ratio = deck_cfg.get("ratio", DEFAULT_RATIO)
    if ratio not in RATIO_PRESETS:
        warn(f"不认识的 ratio={ratio}，回落到 {DEFAULT_RATIO}")
        ratio = DEFAULT_RATIO
    default_w, default_h = RATIO_PRESETS[ratio]

    w_in = float(deck_cfg.get("slide_width_in", default_w))
    h_in = float(deck_cfg.get("slide_height_in", default_h))

    prs = Presentation()
    prs.slide_width = Inches(w_in)
    prs.slide_height = Inches(h_in)
    return prs, w_in, h_in


def add_background(slide, bg_source: Path | io.BytesIO,
                   slide_w_in: float, slide_h_in: float) -> None:
    """背景图占满整页。"""
    if isinstance(bg_source, Path) and not bg_source.exists():
        raise SystemExit(f"背景图不存在: {bg_source}")
    image_input = str(bg_source) if isinstance(bg_source, Path) else bg_source
    slide.shapes.add_picture(
        image_input,
        left=Inches(0), top=Inches(0),
        width=Inches(slide_w_in), height=Inches(slide_h_in),
    )


def add_text_box(slide, tb: dict, slide_w_in: float, slide_h_in: float,
                 default_font: str, default_font_size_pt: float,
                 default_color: str) -> None:
    # 坐标：fraction → inch
    x_in = float(tb["x"]) * slide_w_in
    y_in = float(tb["y"]) * slide_h_in
    w_in = float(tb["w"]) * slide_w_in
    h_in = float(tb["h"]) * slide_h_in

    box = slide.shapes.add_textbox(
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    valign = tb.get("valign", "top")
    tf.vertical_anchor = VALIGN_MAP.get(valign, MSO_ANCHOR.TOP)

    if tb.get("auto_fit", False):
        # python-pptx 的 auto_size 支持有限；这里只设置标志，依靠 PowerPoint 自身做
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    text = tb.get("text", "")
    if not isinstance(text, str):
        text = str(text)

    font_family = tb.get("font_family") or default_font
    if font_family not in SAFE_FONT_SET:
        warn(f"字体 '{font_family}' 不在 SAFE_FONT_SET，跨平台可能掉 fallback")
    font_size_pt = float(tb.get("font_size_pt", default_font_size_pt))
    color_hex = tb.get("color") or default_color
    color_rgb = hex_to_rgb(color_hex)
    bold = bool(tb.get("bold", False))
    italic = bool(tb.get("italic", False))
    align = tb.get("align", "left")
    align_enum = ALIGN_MAP.get(align, PP_ALIGN.LEFT)
    line_spacing = float(tb.get("line_spacing", DEFAULT_LINE_SPACING))

    lines = text.split("\n")
    # 清掉默认空 paragraph
    p0 = tf.paragraphs[0]
    p0.alignment = align_enum
    p0.line_spacing = line_spacing
    run0 = p0.add_run()
    run0.text = lines[0] if lines else ""
    _apply_run_style(run0, font_family, font_size_pt, color_rgb, bold, italic)

    for line in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = align_enum
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        _apply_run_style(r, font_family, font_size_pt, color_rgb, bold, italic)


def _apply_run_style(run, font_family: str, font_size_pt: float,
                     color_rgb: RGBColor, bold: bool, italic: bool) -> None:
    run.font.name = font_family
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = color_rgb
    run.font.bold = bold
    run.font.italic = italic


def render_preview_via_pil(bg_source: Path | io.BytesIO, tbs: list[dict],
                            slide_w_in: float, slide_h_in: float,
                            out_path: Path, dpi: int = 150,
                            default_font: str = DEFAULT_FONT,
                            default_font_size_pt: float = DEFAULT_FONT_SIZE_PT,
                            default_color: str = DEFAULT_COLOR) -> None:
    """
    用 PIL 渲染一张近似预览图（不是真的从 PPTX 导出，只是一个直观对照）。
    PPTX → PNG 真正完美的做法是用 LibreOffice / Keynote / PowerPoint 转，
    这里只做"够用的对照预览"，避免引入额外的重依赖。
    """
    if Image is None:
        warn("没装 Pillow，跳过预览渲染")
        return
    from PIL import Image as PILImage, ImageDraw, ImageFont

    out_path.parent.mkdir(parents=True, exist_ok=True)

    W = int(slide_w_in * dpi)
    H = int(slide_h_in * dpi)
    bg = PILImage.open(bg_source).convert("RGB").resize((W, H), PILImage.LANCZOS)
    draw = ImageDraw.Draw(bg)

    for tb in tbs:
        x = int(float(tb["x"]) * W)
        y = int(float(tb["y"]) * H)
        w = int(float(tb["w"]) * W)
        h = int(float(tb["h"]) * H)

        text = tb.get("text", "")
        font_size_pt = float(tb.get("font_size_pt", default_font_size_pt))
        # 1 pt = dpi/72 px
        font_size_px = int(font_size_pt * dpi / 72)
        color_hex = tb.get("color") or default_color
        color_rgb = _hex_to_tuple(color_hex)

        # PIL 找字体：依次尝试 SAFE_FONT_SET 里的常见路径
        font = _pick_font(tb.get("font_family") or default_font, font_size_px)
        align = tb.get("align", "left")

        # 简易换行
        lines = []
        for raw_line in text.split("\n"):
            lines.extend(_wrap_line(raw_line, font, w, draw))

        line_height = int(font_size_px * float(tb.get("line_spacing", DEFAULT_LINE_SPACING)))
        cur_y = y
        for line in lines:
            tw = _measure_text(draw, line, font)
            if align == "center":
                tx = x + (w - tw) // 2
            elif align == "right":
                tx = x + w - tw
            else:
                tx = x
            draw.text((tx, cur_y), line, fill=color_rgb, font=font)
            cur_y += line_height
            if cur_y > y + h:
                break

    bg.save(out_path)


def _hex_to_tuple(s: str) -> tuple[int, int, int]:
    s = s.strip().lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


_FONT_CACHE: dict[tuple[str, int], Any] = {}


def _pick_font(name: str, size_px: int):
    """从系统常见路径里找一个能用的字体；找不到就用 PIL 默认。"""
    from PIL import ImageFont
    key = (name, size_px)
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]

    candidates = []
    n = name.lower()
    # macOS 字体路径
    if "pingfang" in n:
        candidates += [
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/Hiragino Sans GB.ttc",
        ]
    if "yahei" in n:
        candidates += [
            "C:/Windows/Fonts/msyh.ttc",
            "C:/Windows/Fonts/msyh.ttf",
        ]
    if "arial" in n:
        candidates += [
            "/Library/Fonts/Arial.ttf",
            "C:/Windows/Fonts/arial.ttf",
        ]
    if "helvetica" in n:
        candidates += [
            "/System/Library/Fonts/Helvetica.ttc",
        ]
    # 兜底：找一个 CJK 字体
    candidates += [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "C:/Windows/Fonts/msyh.ttc",
        "/usr/share/fonts/opentype/source-han-sans/SourceHanSansSC-Regular.otf",
    ]

    for path in candidates:
        try:
            f = ImageFont.truetype(path, size_px)
            _FONT_CACHE[key] = f
            return f
        except Exception:
            continue
    # 最后退到默认
    f = ImageFont.load_default()
    _FONT_CACHE[key] = f
    return f


def _measure_text(draw, text: str, font) -> int:
    """跨 PIL 版本兼容的文本宽度测量。"""
    if hasattr(draw, "textlength"):
        return int(draw.textlength(text, font=font))
    bbox = draw.textbbox((0, 0), text, font=font)
    return bbox[2] - bbox[0]


def _wrap_line(line: str, font, max_w: int, draw) -> list[str]:
    """按像素宽度换行；中英文混排做简易切分（按字符切，不按词）。"""
    if not line:
        return [""]
    out = []
    buf = ""
    for ch in line:
        candidate = buf + ch
        if _measure_text(draw, candidate, font) <= max_w:
            buf = candidate
        else:
            if buf:
                out.append(buf)
            buf = ch
    if buf:
        out.append(buf)
    return out or [""]


# ─────────────────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────────────────

def parse_args() -> argparse.Namespace:
    ap = argparse.ArgumentParser(
        description=__doc__,
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    ap.add_argument("input", type=Path, help="layout.json (单页) 或 deck.json (多页)")
    ap.add_argument("-o", "--out", required=True, type=Path, help="输出 .pptx 路径")
    ap.add_argument("--preview-dir", type=Path,
                    help="可选：每页同时生成 PIL 近似预览图到此目录（用于和 HTML 编辑器视觉对照）")
    ap.add_argument("--preview-dpi", type=int, default=150,
                    help="预览图 DPI，默认 150")
    return ap.parse_args()


def main() -> None:
    args = parse_args()
    deck_cfg, slides, base_dir = load_deck(args.input.resolve())

    prs, slide_w_in, slide_h_in = make_presentation(deck_cfg)
    default_font = deck_cfg.get("default_font", DEFAULT_FONT)
    default_font_size_pt = float(deck_cfg.get("default_font_size_pt", DEFAULT_FONT_SIZE_PT))
    default_color = deck_cfg.get("default_color", DEFAULT_COLOR)

    blank_layout = prs.slide_layouts[6]  # 空白版式

    for i, s in enumerate(slides, start=1):
        sid = s.get("id") or f"{i:02d}"
        bg_rel = s.get("background")
        if not bg_rel:
            raise SystemExit(f"slide {sid} 缺少 background 字段")
        bg_source = resolve_background_source(bg_rel, base_dir)

        slide = prs.slides.add_slide(blank_layout)
        add_background(slide, clone_background_source(bg_source), slide_w_in, slide_h_in)

        for tb in s.get("text_boxes", []):
            try:
                add_text_box(slide, tb, slide_w_in, slide_h_in,
                             default_font, default_font_size_pt, default_color)
            except Exception as e:
                warn(f"slide {sid} 文字框写入失败: {e} (tb={tb.get('id', '?')})")

        if args.preview_dir:
            preview_out = args.preview_dir.resolve() / f"slide_{i:02d}.png"
            try:
                render_preview_via_pil(
                    clone_background_source(bg_source), s.get("text_boxes", []),
                    slide_w_in, slide_h_in, preview_out,
                    dpi=args.preview_dpi,
                    default_font=default_font,
                    default_font_size_pt=default_font_size_pt,
                    default_color=default_color,
                )
            except Exception as e:
                warn(f"slide {sid} 预览图渲染失败: {e}")

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.out))
    print(f"[ok] {len(slides)} 页 → {args.out}")
    if args.preview_dir:
        print(f"[ok] 预览图 → {args.preview_dir}")


if __name__ == "__main__":
    main()
